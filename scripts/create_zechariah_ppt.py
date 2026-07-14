#!/usr/bin/env python3
"""Generate Zechariah 2:8-9 presentation from document content."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

OUTPUT_PATH = "/workspace/Zechariah_2_8-9_Presentation.pptx"

# Colors
TITLE_COLOR = RGBColor(0x1A, 0x3A, 0x5C)
ACCENT_COLOR = RGBColor(0x8B, 0x45, 0x13)
BODY_COLOR = RGBColor(0x33, 0x33, 0x33)
VERSE_COLOR = RGBColor(0x2E, 0x5A, 0x88)


def set_slide_background(slide, r=245, g=248, b=252):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


def add_title_slide(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_background(slide)
    box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(1.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        box2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1))
        tf2 = box2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = ACCENT_COLOR
        p2.alignment = PP_ALIGN.CENTER
    return slide


def add_section_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, 26, 58, 92)
    box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    return slide


def add_content_slide(prs, title, bullets, font_size=16):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    # Title
    tbox = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(9.2), Inches(0.8))
    tp = tbox.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(24)
    tp.font.bold = True
    tp.font.color.rgb = TITLE_COLOR

    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(5.8))
    tf = body.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            text, level, bold, color = item[0], item[1] if len(item) > 1 else 0, item[2] if len(item) > 2 else False, item[3] if len(item) > 3 else BODY_COLOR
        else:
            text, level, bold, color = item, 0, False, BODY_COLOR
        p.text = text
        p.level = level
        p.font.size = Pt(font_size - level * 2)
        p.font.bold = bold
        p.font.color.rgb = color
        p.space_after = Pt(6)
    return slide


def add_verse_slide(prs, title, english_lines, tamil_lines=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    tbox = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(9.2), Inches(0.7))
    tp = tbox.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(22)
    tp.font.bold = True
    tp.font.color.rgb = TITLE_COLOR

    y = 1.0
    if english_lines:
        ebox = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(2.5))
        etf = ebox.text_frame
        etf.word_wrap = True
        for i, line in enumerate(english_lines):
            p = etf.paragraphs[0] if i == 0 else etf.add_paragraph()
            p.text = line
            p.font.size = Pt(15)
            p.font.italic = True
            p.font.color.rgb = VERSE_COLOR
            p.space_after = Pt(8)
        y += 2.6

    if tamil_lines:
        tbox2 = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(3))
        ttf = tbox2.text_frame
        ttf.word_wrap = True
        for i, line in enumerate(tamil_lines):
            p = ttf.paragraphs[0] if i == 0 else ttf.add_paragraph()
            p.text = line
            p.font.size = Pt(14)
            p.font.color.rgb = ACCENT_COLOR
            p.space_after = Pt(6)
    return slide


def add_text_slide(prs, title, paragraphs, font_size=15):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    tbox = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(9.2), Inches(0.7))
    tp = tbox.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(22)
    tp.font.bold = True
    tp.font.color.rgb = TITLE_COLOR

    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(5.8))
    tf = body.text_frame
    tf.word_wrap = True
    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = para
        p.font.size = Pt(font_size)
        p.font.color.rgb = BODY_COLOR
        p.space_after = Pt(10)
    return slide


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. Title
    add_title_slide(
        prs,
        "Zechariah 2:8-9",
        "Connecting to the First War in Human History — Battle of Siddim",
    )

    # 2. Main verses
    add_verse_slide(
        prs,
        "Zechariah 2:8-9",
        [
            "8 For this is what the Lord Almighty says: \"After the Glorious One has sent me against the nations that have plundered you—for whoever touches you touches the apple of his eye—",
            "9 I will surely raise my hand against them so that their slaves will plunder them.[b] Then you will know that the Lord Almighty has sent me.\"",
        ],
        [
            "சகரியா 2",
            "8 பிற்பாடு மகிமையுண்டாகுமென்று சேனைகளின் கர்த்தர் சொல்லுகிறார்; உங்களைக் கொள்ளையிட்ட ஜாதிகளிடத்துக்கு என்னை அனுப்பினார்; உங்களைத் தொடுகிறவன் அவருடைய கண்மணியைத் தொடுகிறான்.",
            "9 இதோ, நான் என் கையை அவர்களுக்கு விரோதமாக அசைப்பேன்; அதினால் அவர்கள் தங்கள் அடிமைகளுக்குக் கொள்ளையாவார்கள்; அப்பொழுது சேனைகளின் கர்த்தர் என்னை அனுப்பினாரென்று அறிவீர்கள்.",
        ],
    )

    # 3. Section - Battle of Siddim
    add_section_slide(prs, "Connecting to the First War in Human being History\nBattle of Siddim")

    # 4. Side 1 intro
    add_content_slide(
        prs,
        "SIDE 1 — The 4 Kings (The Aggressors / Eastern Coalition)",
        [
            "These were large, established empires from Mesopotamia:",
            ("1. Chedorlaomer — King of Elam (present-day southwest Iran)", 0, True, TITLE_COLOR),
            "The leader of the coalition. His name in Elamite is Kudur-Lagomer. Elam was the dominant superpower of the era — its empire stretched from the Caspian Sea in the north to the Persian Gulf in the south, and its diplomatic and military reach extended all the way west into Syria and the Levant. This was the mightiest empire of the ancient Near East at that time.",
        ],
        font_size=14,
    )

    # 5. Kings 2-4
    add_content_slide(
        prs,
        "SIDE 1 — The 4 Kings (continued)",
        [
            ("2. Amraphel — King of Shinar (Sumer / Babylon)", 0, True, TITLE_COLOR),
            '"Shinar" is the Bible\'s name for Sumer, the civilisation of Mesopotamia that produced Babylon. The article makes a fascinating connection — Amraphel may be identified with Hammurabi, the famous Babylonian king. The Semitic form of his name (Amurapi-ili) closely matches the Hebrew Amraphel. He was at this stage still a relatively minor regional ruler who called the king of Elam his "father."',
            ("3. Arioch — King of Ellasar (Larsa)", 0, True, TITLE_COLOR),
            "Ellasar is identified as the ancient kingdom of Larsa, which existed only for about 150 years during the 19th century BC. Arioch's name matches the last known king of Larsa — Rim-Sin I, whose Semitic name is Eri-aku. The article argues this is a strong historical identification.",
            ("4. Tidal — King of Goiim (Hittites / \"Nations\")", 0, True, TITLE_COLOR),
            'Goiim simply means "nations" or "peoples" in Hebrew. Tidal is a recognised Hittite name — Tudhalia — and the Hittite empire in modern-day Turkey was a confederation of many different tribes and peoples, which is why its king was called "king of nations." A mysterious proto-Tudhalia figure from the 19th–17th centuries BC is the likely match.',
        ],
        font_size=13,
    )

    # 6. Side 2
    add_content_slide(
        prs,
        "SIDE 2 — The 5 Kings (The Defenders / Jordan Valley City-States)",
        [
            "These were small, weak city-states in the region of the Dead Sea that had been paying tribute to Chedorlaomer for 13 years:",
            "King\tCity",
            "Bera\tSodom",
            "Birsha\tGomorrah",
            "Shinab\tAdmah",
            "Shemeber\tZeboiim",
            "Unnamed\tBela / Zoar",
            "",
            "In the 13th year of their oppression, these five kings rebelled and refused to pay tribute — which triggered the entire invasion.",
        ],
        font_size=15,
    )

    # 7. March through Levant
    add_content_slide(
        prs,
        "What Happened Before the Main Battle — The March Through the Levant",
        [
            "Before even reaching the Jordan Valley, Chedorlaomer's coalition swept through the entire Levant in a devastating campaign, attacking multiple peoples along the way. According to Genesis 14:5-7, they destroyed:",
            "• The Rephaim at Ashteroth-karnaim",
            "• The Zuzim at Ham",
            "• The Emim in the plain of Kiriathaim",
            "• The Horites in the hill country of Seir",
            "• The Amalekites",
            "• The western Amorites",
            "",
            "This was a systematic military campaign of maximum destruction across the entire region before they even engaged the five kings.",
        ],
        font_size=14,
    )

    # 8. Giants
    add_content_slide(
        prs,
        "Who Are the Giants? — The Rephaim and Emim",
        [
            "This is crucial and the article touches on it. The peoples destroyed in Genesis 14:5-6 were no ordinary tribes:",
            ("The Rephaim", 0, True, TITLE_COLOR),
            'The Hebrew word Rephaim refers to a race of giants. They appear throughout the Old Testament as a fearsome giant people. In Deuteronomy 2:11, the Emim are described as "a people great and tall as the Anakim" — the same Anakim that terrified the 12 Israelite spies centuries later (Numbers 13:33). The Rephaim are also connected to Og, king of Bashan, whose bed was 9 cubits long (roughly 4 metres). They were a warrior race of enormous physical stature.',
            ("The Emim", 0, True, TITLE_COLOR),
            'Their name literally means "the terrifying ones" or "the fearsome ones." Deuteronomy 2:10-11 describes them as tall as the Anakim, which is the same category of giant people. They were the original inhabitants of the land of Moab before being displaced.',
        ],
        font_size=13,
    )

    # 9. Giants continued
    add_content_slide(
        prs,
        "Who Are the Giants? (continued)",
        [
            ("The Zuzim (also called Zamzummim)", 0, True, TITLE_COLOR),
            'Another giant people group, also described in Deuteronomy 2:20 as "a people great and numerous and tall as the Anakim."',
            ("The Horites", 0, True, TITLE_COLOR),
            "Inhabitants of the mountain of Seir, later displaced by the descendants of Esau.",
            "",
            "So when Chedorlaomer and his coalition swept through the Levant on their way to the Jordan Valley, they were not just crushing small villages — they were destroying entire nations of giant warrior peoples that no one else had been able to conquer. This makes their military power even more staggering, and makes Abraham's eventual defeat of them even more extraordinary.",
        ],
        font_size=14,
    )

    # 10. Summary
    add_text_slide(
        prs,
        "Summary in One Paragraph",
        [
            "Four eastern superpowers — led by the Elamite empire under Chedorlaomer — crushed five small Jordan Valley city-states that had rebelled after 13 years of paying tribute. On the way, they destroyed entire nations of giant warrior peoples including the Rephaim, Emim and Zuzim. They plundered Sodom and Gomorrah and took captives including Lot. Abraham then pursued them 200 kilometres with 318 household servants, launched a night guerrilla attack with divided forces, and utterly defeated all four armies. Every captive was freed. The victory was so complete that it may have triggered the collapse of the Elamite empire and the rise of Babylon — making a private rescue mission by one man the event that reshaped the entire ancient world.",
        ],
        font_size=14,
    )

    # 11. Wider significance
    add_text_slide(
        prs,
        "The Wider Significance the Article Highlights",
        [
            "The Armstrong Institute article makes a remarkable historical argument about what happened after the battle. Shortly after this period in history, the Elamite empire collapsed. Larsa was conquered. The Hittite world was destabilised. And the Amorites — the very people among whom Abraham lived and, in whose territory, Chedorlaomer had attacked — rose to become the new dominant power of Mesopotamia under Hammurabi.",
            "The author asks whether Abraham's humiliation of Elam in the Levant may have been the event that tipped the balance — weakening Chedorlaomer's prestige and authority enough that Hammurabi and the Amorites were able to overthrow Elam and Larsa and establish the Babylonian empire.",
        ],
        font_size=14,
    )

    # 12. Five Deep Links section
    add_section_slide(prs, "Five Deep Links and Promises\nmade by God in Zechariah 2:8-9")

    # 13. Link 1
    add_verse_slide(
        prs,
        "Link 1 — The Same Hand",
        [
            'Genesis 14 — God delivered the enemies into Abraham\'s hand. Melchizedek specifically used the word hand: "delivered your enemies into your hand."',
            'Zechariah 2:9 — God says "I will shake my hand over them." His hand moves — and the result is the enemy becomes plunder.',
        ],
        [
            "Tamil – Genesis 14:20 (20 உன் சத்துருக்களை உன் கையில் ஒப்புக்கொடுத்த உன்னதமான தேவனுக்கு ஸ்தோத்திரம் என்று சொன்னான். இவனுக்கு ஆபிராம் எல்லாவற்றிலும் தசமபாகம் கொடுத்தான்.)",
            "9 இதோ, நான் என் கையை அவர்களுக்கு விரோதமாக அசைப்பேன்",
        ],
    )

    add_text_slide(
        prs,
        "Link 1 — The Same Hand (continued)",
        [
            "In both texts, the hand is the instrument of reversal. In Genesis 14, God's hand works through Abraham's 318 men. In Zechariah 2:9, God shakes His own hand directly over the nations.",
            "The progression is significant. In Genesis 14, God works through a human hand. In Zechariah 2:9, God acts with His own hand.",
        ],
        font_size=16,
    )

    # 14. Link 2 (NO Testimonial: Jefferson vs Balaji)
    add_content_slide(
        prs,
        "Link 2 — The Plunderer Becomes the Plundered",
        [
            "This is the most direct verbal connection between the two passages.",
            "Genesis 14 — Chedorlaomer's coalition plundered Sodom and took Lot captive. They were the powerful ones, the captors, the plunderers. Then Abraham came and reversed everything completely. The plunderers lost everything they had taken.",
            'Zechariah 2:9 — "They shall become plunder for those who served them."',
            "The ones who plundered Israel — Babylon, Persia, the nations — will themselves become plunder for those they once enslaved.",
            "This is not just a coincidence of language. It is the same divine principle operating at a larger scale.",
            "Abraham demonstrated in Genesis 14 what God declared He would do permanently in Zechariah 2:9. The small battle in the Valley of Siddim was a preview — a type — of God's permanent policy toward those who touch His people.",
        ],
        font_size=14,
    )

    # 15. Link 3 (NO Testimony: Jefferson vs Wein Chein)
    add_verse_slide(
        prs,
        "Link 3 — The Apple of His Eye Principle",
        [
            'Zechariah 2:8 — the verse immediately before your text — says:',
            '"For thus said the Lord of hosts, after his glory sent me to the nations who plundered you, for he who touches you touches the apple of his eye."',
        ],
        [
            "உங்களைத் தொடுகிறவன் அவருடைய கண்மணியைத் தொடுகிறான்",
        ],
    )

    add_text_slide(
        prs,
        "Link 3 — The Apple of His Eye (continued)",
        [
            "God was already operating the Zechariah 2:8-9 principle in Genesis 14 — He just had not put it into words yet.",
            "Every time a nation touched Abraham's household, they touched the apple of God's eye. God's hand shook — and they became plunder.",
        ],
        font_size=16,
    )

    # 16. Link 4
    add_verse_slide(
        prs,
        'Link 4 — "You Shall Know" — The Authenticating Moment',
        [
            'Zechariah 2:9 ends with: "Then you will know that the Lord of hosts has sent me."',
            "Every time God fights, He is also revealing Himself. The battle is not just military — it is theological.",
        ],
        [
            "அப்பொழுது சேனைகளின் கர்த்தர் என்னை அனுப்பினாரென்று அறிவீர்கள்",
        ],
    )

    # 17. What We Must Do section
    add_section_slide(prs, "What We Must Do For God\nto Lead the Battle")

    # 18. Exodus 14:14
    add_verse_slide(
        prs,
        "Exodus 14:14",
        [
            '"The Lord will fight for you — you need only to be still."',
        ],
        [
            "யாத்திராகமம் 14:14",
            "கர்த்தர் உங்களுக்காக யுத்தம் பண்ணுவார்; நீங்கள் சும்மாயிருப்பீர்கள் என்றான்.",
        ],
    )

    add_text_slide(
        prs,
        "What we must do for the Zech 2:8-9 to be fulfilled in our life",
        [
            "Be Still means not to relay on yourselves, Not to Trust you but put full trust in God.",
        ],
        font_size=16,
    )

    # 19. Point 1
    add_content_slide(
        prs,
        "1. Acknowledge You Cannot Win Without Him",
        [
            "அவர் இல்லாமல் உங்களால் வெற்றி பெற முடியாது என்பதை ஒப்புக்கொள்ளுங்கள்.",
            "",
            "The First condition is always hard for the Strong people. The people who has more self confidence on themselves, Because these people feel that they can face any strong challenge without the need of others help (God's help)",
            "Gideon's 300 — God reduced the army deliberately so no one could say \"my own strength saved me.\"",
            "Abharam 318 – God helped Abharam to destroy a skilled warrior and his allies",
            'Reality / Truth: The moment you say "I can handle this" — God steps back. The moment you say "I cannot handle this without You" — He steps forward.',
            "Practical Life:",
            "Before any battle — whether spiritual, financial, relational or physical — come to God with empty hands. Not with your plan for Him to bless. With your helplessness for Him to fill.",
        ],
        font_size=13,
    )

    # 20. Point 2
    add_verse_slide(
        prs,
        "2. Seek God First — Before You Make a Single Move",
        [
            "God works in different ways so always ask the plan – Email draft should be God's words",
            "This is the condition that Israel broke most often and paid for most severely.",
            "",
            '2 Samuel 5:19',
            'so David inquired of the Lord, "Shall I go and attack the Philistines? Will you deliver them into my hands?"',
            'The Lord answered him, "Go, for I will surely deliver the Philistines into your hands."',
            "",
            '2 Samuel 5:23',
            'David inquired of the Lord, and he answered, "Do not go straight up, but circle around behind them and attack them in front of the poplar trees."',
        ],
        [
            "பெலிஸ்தருக்கு விரோதமாய்ப்போகலாமா, அவர்களை என் கையில் ஒப்புக்கொடுப்பீரா என்று தாவீது கர்த்தரிடத்தில் விசாரித்தபோது, கர்த்தர்: போ, பெலிஸ்தரை உன் கையில் நிச்சயமாய் ஒப்புக்கொடுப்பேன் என்று தாவீதுக்குச் சொன்னார்.",
            "23 தாவீது கர்த்தரிடத்தில் விசாரித்ததற்கு, அவர்: நீ நேராய்ப் போகாமல், அவர்களுக்குப் பின்னாலே சுற்றி, முசுக்கட்டைச் செடிகளுக்கு எதிரேயிருந்து, அவர்கள்மேல் பாய்ந்து,",
        ],
    )

    add_text_slide(
        prs,
        "2. Seek God First (continued)",
        [
            "Different strategy, same enemy, same God — because David asked again instead of assuming.",
            "Practical application:",
            "Do not carry yesterday's answer into today's battle. Seek God fresh every time. What He said last year may not be what He is saying today. His strategy changes — His faithfulness does not.",
        ],
        font_size=15,
    )

    # 21. Point 3 (NO Testimony: Jefferson Leg issue)
    add_verse_slide(
        prs,
        "3. Obey Exactly — Not Approximately",
        [
            "Naman Obeyed 100%!",
            "",
            "2 Kings 5:14:",
            "14 So he went down and dipped himself in the Jordan seven times, as the man of God had told him, and his flesh was restored and became clean like that of a young boy.",
        ],
        [
            "14 அப்பொழுது அவன் இறங்கி, தேவனுடைய மனுஷன் வார்த்தையின்படியே யோர்தானில் ஏழுதரம் முழுகினபோது, அவன் மாம்சம் ஒரு சிறுபிள்ளையின் மாம்சத்தைப்போல மாறி, அவன் சுத்தமானான்.",
        ],
    )

    add_content_slide(
        prs,
        "3. Obey Exactly — Not Approximately (continued)",
        [
            "Naaman the leper — told to dip in the Jordan seven times. He dipped six times and came up still leprous. The seventh dip — full obedience — produced full healing.",
            "",
            "Dis-obeyed:",
            "Saul — told to completely destroy the Amalekites. He destroyed most of them and kept the king and best animals for what seemed like a good reason. God called it rebellion.",
            "",
            "The partial obedience of Saul cost him his kingdom. The complete obedience of Joshua gave him a nation.",
        ],
        font_size=14,
    )

    # 22. Point 4
    add_content_slide(
        prs,
        "4. Remove Sin From the Heart",
        [
            ("The Sin of Achan:", 0, True, TITLE_COLOR),
            ("The Order of Events", 0, True, TITLE_COLOR),
            "1. The Sin (Secret): Achan stole the forbidden items during the fall of Jericho and hid them under his tent (Joshua 6:27-7:1).",
            "2. The Battle (Public): Joshua, completely unaware of Achan's theft, sent the small army of 3,000 men to attack Ai (Joshua 7:2-4).",
            "3. The Defeat: The small army was routed, and 36 Israelite soldiers were killed (Joshua 7:4-5).",
            "4. The Revelation: Joshua wept and prayed to God, and only then did God reveal to him that someone in the camp had sinned (Joshua 7:10-11)",
            "",
            'God did not say "Achan has sinned." He said "Israel has sinned."',
        ],
        font_size=13,
    )

    add_text_slide(
        prs,
        "4. Remove Sin From the Heart (continued)",
        [
            "The Sin in our heart will not complete the Promised Verse for us:",
            "Deal with these first. Do not go into battle carrying what will stop God's hand.",
        ],
        font_size=16,
    )

    # 23. Point 5
    add_text_slide(
        prs,
        "5. Praise God Before Not After (worship is the weapon)",
        [
            "Before the difficult meeting, before the medical procedure, before the financial crisis hearing, before the confrontation — worship first. Not as a formula but as a declaration that God is already Lord over what you are about to face.",
            "We are declaring and stamping a statement that God is in Control",
        ],
        font_size=15,
    )

    # 24. Point 6
    add_content_slide(
        prs,
        "6. Give God the Glory — Before, During and After",
        [
            "After winning the battle of Genesis 14, the king of Sodom offered Abraham all the recovered goods. Abraham refused everything — not even a sandal strap. He said:",
            '"I have lifted my hand to the Lord, God Most High, that I would not take a thread or sandal strap lest you should say \'I have made Abraham rich.\'"',
            "He protected the glory. He would not allow the king of Sodom to even partially claim credit for what God had done.",
            "",
            'When God gives you the victory, be careful how you speak about it. Do not say "I worked hard and it paid off." Do not say "I was wise enough to make the right decisions." Acknowledge God specifically and completely. The moment you absorb the glory that belongs to Him, you are positioning yourself for the next battle without His protection.',
        ],
        font_size=13,
    )

    # 25. Conclusion
    add_verse_slide(
        prs,
        "Conclusion",
        [
            "2 Chronicles 16:9",
            '"For the eyes of the Lord range throughout the earth to strengthen those whose hearts are fully committed to him."',
        ],
        [
            "தம்மைப்பற்றி உத்தம இருதயத்தோடிருக்கிறவர்களுக்குத் தம்முடைய வல்லமையை விளங்கப்பண்ணும்படி, கர்த்தருடைய கண்கள் பூமியெங்கும் உலாவிக்கொண்டிருக்கிறது;",
        ],
    )

    add_text_slide(
        prs,
        "Conclusion (continued)",
        [
            "He is looking. Right now. Not for the most talented, the most experienced, the most resourced. He is looking for the most surrendered — the one whose heart is fully His.",
            "That is the condition. That is always the condition. And when He finds it — He fights.",
        ],
        font_size=18,
    )

    prs.save(OUTPUT_PATH)
    print(f"Saved: {OUTPUT_PATH}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_presentation()
